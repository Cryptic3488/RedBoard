"""
RedBoard presentation builder — python-pptx
Generates presentation/RedBoard.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os
from pathlib import Path

# ── Paths ──────────────────────────────────────────────────────────────────────
BASE = Path(__file__).parent.parent
SHOTS = BASE / "presentation" / "screenshots"
OUT   = BASE / "presentation" / "RedBoard.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# ── Brand colors ───────────────────────────────────────────────────────────────
RED    = RGBColor(0xC8, 0x10, 0x2E)   # Denison brand red
CREAM  = RGBColor(0xF5, 0xF0, 0xEB)   # light mode bg
DARK   = RGBColor(0x1C, 0x1C, 0x1E)   # dark background
CARD   = RGBColor(0x2C, 0x2C, 0x2E)   # card surface
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xA0, 0xA0, 0xA8)   # muted label

# Slide size: widescreen 16:9
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ────────────────────────────────────────────────────────────────────

def fill_bg(slide, color: RGBColor):
    """Solid background fill for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_image(slide, path, x, y, w=None, h=None):
    if not Path(path).exists():
        return None
    if w and h:
        return slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    elif w:
        return slide.shapes.add_picture(str(path), x, y, width=w)
    elif h:
        return slide.shapes.add_picture(str(path), x, y, height=h)
    else:
        return slide.shapes.add_picture(str(path), x, y)


def red_bar(slide, width=Inches(0.08)):
    """Left red accent bar (full height)."""
    add_rect(slide, 0, 0, width, H, RED)


def slide_label(slide, text, y=Inches(0.28)):
    add_text(slide, text.upper(), Inches(0.22), y, Inches(4), Inches(0.35),
             size=9, color=LGRAY, bold=True)


def heading(slide, text, y=Inches(0.55), size=38, color=WHITE):
    add_text(slide, text, Inches(0.22), y, Inches(12.8), Inches(1.2),
             size=size, bold=True, color=color)


def sub(slide, text, y=Inches(1.45), size=18, color=LGRAY):
    add_text(slide, text, Inches(0.22), y, Inches(12.8), Inches(0.6),
             size=size, color=color)


def bullet_block(slide, items, x, y, w, item_size=16, color=WHITE, gap=Inches(0.02)):
    """Render a list of bullet strings into stacked text boxes."""
    row_h = Inches(0.38)
    for i, item in enumerate(items):
        add_text(slide, f"• {item}", x, y + i * (row_h + gap), w, row_h,
                 size=item_size, color=color)


# ── Slide builders ─────────────────────────────────────────────────────────────

def slide_cover():
    sl = prs.slides.add_slide(BLANK)
    fill_bg(sl, DARK)

    # Red accent bar (wide)
    add_rect(sl, 0, 0, Inches(0.5), H, RED)

    # Logo placeholder — "D" in brand red as large letter
    add_text(sl, "D", Inches(0.8), Inches(0.6), Inches(1.5), Inches(1.5),
             size=96, bold=True, color=RED)

    # Title
    add_text(sl, "RedBoard", Inches(0.8), Inches(1.9), Inches(10), Inches(1.4),
             size=64, bold=True, color=WHITE)

    # Subtitle
    add_text(sl, "Denison University Women's Basketball",
             Inches(0.8), Inches(3.1), Inches(10), Inches(0.7),
             size=24, color=CREAM)

    # Tagline
    add_text(sl, "One platform. Coach content. Player focus.",
             Inches(0.8), Inches(3.85), Inches(10), Inches(0.6),
             size=18, color=LGRAY, italic=True)

    # Bottom strip
    add_rect(sl, 0, H - Inches(0.5), W, Inches(0.5), RED)
    add_text(sl, "CONFIDENTIAL — INTERNAL USE ONLY",
             Inches(0.2), H - Inches(0.42), W - Inches(0.4), Inches(0.35),
             size=11, color=WHITE, align=PP_ALIGN.CENTER)


def slide_problem():
    sl = prs.slides.add_slide(BLANK)
    fill_bg(sl, DARK)
    red_bar(sl)
    slide_label(sl, "The Problem")
    heading(sl, "Coaching ran on group chats,\nshared drives, and printed binders.")

    items = [
        "Film buried in Hudl links scattered across text threads",
        "Stat sheets emailed as Excel attachments — no version control",
        "Wellness check-ins done verbally or forgotten entirely",
        "Playbook PDFs printed, lost, or outdated",
        "No single source of truth for players or staff",
    ]
    bullet_block(sl, items, Inches(0.4), Inches(2.1), Inches(12.5),
                 item_size=18, color=CREAM)


def slide_solution():
    sl = prs.slides.add_slide(BLANK)
    fill_bg(sl, DARK)
    red_bar(sl)
    slide_label(sl, "The Solution")
    heading(sl, "RedBoard replaces the chaos\nwith one mobile-first platform.")

    boxes = [
        ("Film", "Hudl clips & video\nshared instantly"),
        ("Stats", "CSV upload → published\nwith trend analysis"),
        ("Wellness", "Daily check-in forms\nbuild player data"),
        ("Playbook", "Virtual binder with\nfolders & files"),
    ]

    box_w = Inches(2.8)
    box_h = Inches(2.2)
    gap   = Inches(0.25)
    start_x = Inches(0.4)
    y = Inches(2.0)

    for i, (title, desc) in enumerate(boxes):
        x = start_x + i * (box_w + gap)
        add_rect(sl, x, y, box_w, box_h, CARD)
        add_rect(sl, x, y, box_w, Inches(0.07), RED)
        add_text(sl, title, x + Inches(0.18), y + Inches(0.2), box_w - Inches(0.3),
                 Inches(0.5), size=20, bold=True, color=RED)
        add_text(sl, desc, x + Inches(0.18), y + Inches(0.75), box_w - Inches(0.3),
                 Inches(1.2), size=14, color=CREAM)


def slide_stack():
    sl = prs.slides.add_slide(BLANK)
    fill_bg(sl, DARK)
    red_bar(sl)
    slide_label(sl, "Technology")
    heading(sl, "Built on a modern,\nserverless stack.", size=34)

    rows = [
        ("Frontend",   "React 18 (TypeScript) · Vite 5 · Tailwind CSS 3"),
        ("Routing",    "React Router v6"),
        ("Backend",    "Supabase — Auth, PostgREST, Realtime, Storage"),
        ("Database",   "PostgreSQL with Row Level Security (RLS)"),
        ("File parse", "SheetJS (xlsx) — client-side CSV/XLSX parsing"),
        ("Deploy",     "Static Vite build → any CDN or hosting platform"),
    ]

    row_h = Inches(0.68)
    start_y = Inches(2.0)

    for i, (label, value) in enumerate(rows):
        y = start_y + i * row_h
        bg = CARD if i % 2 == 0 else RGBColor(0x24, 0x24, 0x26)
        add_rect(sl, Inches(0.22), y, Inches(12.8), row_h - Inches(0.05), bg)
        add_text(sl, label, Inches(0.4), y + Inches(0.12), Inches(2.2),
                 Inches(0.45), size=13, bold=True, color=RED)
        add_text(sl, value, Inches(2.8), y + Inches(0.12), Inches(10.1),
                 Inches(0.45), size=14, color=CREAM)


def slide_architecture():
    sl = prs.slides.add_slide(BLANK)
    fill_bg(sl, DARK)
    red_bar(sl)
    slide_label(sl, "Architecture")
    heading(sl, "No custom server.\nAll security lives in the database.", size=32)

    layers = [
        ("Browser (React App)",             "supabase-js · anon key · HTTPS",          RGBColor(0x1E, 0x3A, 0x5F)),
        ("Supabase Platform",               "Auth (JWT)  ·  PostgREST  ·  Storage  ·  Realtime", RGBColor(0x1A, 0x4A, 0x2E)),
        ("PostgreSQL + RLS",                "Row Level Security enforces all access rules",        RGBColor(0x4A, 0x1A, 0x1A)),
    ]

    box_h = Inches(1.1)
    gap   = Inches(0.22)
    y0    = Inches(2.05)
    bx    = Inches(0.4)
    bw    = Inches(12.5)

    for i, (title, desc, color) in enumerate(layers):
        y = y0 + i * (box_h + gap)
        add_rect(sl, bx, y, bw, box_h, color)
        add_rect(sl, bx, y, Inches(0.07), box_h, RED)
        add_text(sl, title, bx + Inches(0.22), y + Inches(0.12), bw,
                 Inches(0.45), size=18, bold=True, color=WHITE)
        add_text(sl, desc, bx + Inches(0.22), y + Inches(0.58), bw,
                 Inches(0.38), size=13, color=LGRAY)

        if i < len(layers) - 1:
            arr_x = bx + bw / 2
            arr_y = y + box_h + Inches(0.04)
            add_text(sl, "▼", arr_x - Inches(0.15), arr_y, Inches(0.3),
                     Inches(0.2), size=14, color=RED, align=PP_ALIGN.CENTER)

    note = ("Every query carries the user's JWT. "
            "auth.uid() is evaluated by Postgres at query time — "
            "players can never read another player's data.")
    add_text(sl, note, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.6),
             size=13, color=LGRAY, italic=True)


def slide_feature_mobile(label, title, body_lines, img_name):
    """Two-column layout: text left, phone screenshot right."""
    sl = prs.slides.add_slide(BLANK)
    fill_bg(sl, DARK)
    red_bar(sl)
    slide_label(sl, label)
    add_text(sl, title, Inches(0.22), Inches(0.52), Inches(7.8), Inches(1.3),
             size=30, bold=True, color=WHITE)

    bullet_block(sl, body_lines, Inches(0.4), Inches(2.0), Inches(7.4),
                 item_size=15, color=CREAM)

    img_path = SHOTS / img_name
    if img_path.exists():
        # Phone screenshot — tall and narrow
        add_image(sl, img_path, Inches(8.6), Inches(0.25), h=Inches(6.9))


def slide_feature_desktop(label, title, body_lines, img_name):
    """Full-width screenshot with text overlay panel at bottom."""
    sl = prs.slides.add_slide(BLANK)
    fill_bg(sl, DARK)
    red_bar(sl)

    img_path = SHOTS / img_name
    if img_path.exists():
        add_image(sl, img_path, Inches(0.08), Inches(0.0), w=W - Inches(0.08))

    # Darkened overlay strip at bottom
    overlay_h = Inches(2.5)
    add_rect(sl, 0, H - overlay_h, W, overlay_h, RGBColor(0x0A, 0x0A, 0x0C))

    slide_label(sl, label, y=H - overlay_h + Inches(0.15))
    add_text(sl, title, Inches(0.22), H - overlay_h + Inches(0.42), Inches(12.8),
             Inches(0.65), size=22, bold=True, color=WHITE)

    bullet_text = "   ·   ".join(body_lines)
    add_text(sl, bullet_text, Inches(0.22), H - overlay_h + Inches(1.1),
             Inches(12.8), Inches(0.9), size=13, color=CREAM)


def slide_security():
    sl = prs.slides.add_slide(BLANK)
    fill_bg(sl, DARK)
    red_bar(sl)
    slide_label(sl, "Security")
    heading(sl, "Default-deny RLS.\nEvery row protected at the database layer.", size=30)

    rows = [
        ("profiles",              "Own row only",                      "All rows"),
        ("film_posts",            "Posts targeted to them",            "All posts"),
        ("stat_uploads/entries",  "Published uploads · own entries",   "All uploads & entries"),
        ("wellness_forms",        "Active form only",                  "All forms"),
        ("wellness_responses",    "Own responses",                     "All responses"),
        ("playbook_folders/files","All (read-only)",                   "Full CRUD"),
    ]

    hdr_y  = Inches(2.0)
    row_h  = Inches(0.58)
    col_w  = [Inches(4.2), Inches(4.1), Inches(4.0)]
    col_x  = [Inches(0.4), Inches(4.65), Inches(8.8)]
    hdrs   = ["Table", "Player Access", "Admin Access"]

    # Header row
    for j, (hdr, cx, cw) in enumerate(zip(hdrs, col_x, col_w)):
        add_rect(sl, cx, hdr_y, cw - Inches(0.05), row_h, RED)
        add_text(sl, hdr, cx + Inches(0.12), hdr_y + Inches(0.1),
                 cw, row_h - Inches(0.1), size=14, bold=True, color=WHITE)

    for i, (tbl, player, admin) in enumerate(rows):
        y    = hdr_y + (i + 1) * row_h
        vals = [tbl, player, admin]
        bg   = CARD if i % 2 == 0 else RGBColor(0x24, 0x24, 0x26)
        for j, (val, cx, cw) in enumerate(zip(vals, col_x, col_w)):
            add_rect(sl, cx, y, cw - Inches(0.05), row_h - Inches(0.04), bg)
            add_text(sl, val, cx + Inches(0.12), y + Inches(0.1),
                     cw - Inches(0.2), row_h - Inches(0.15), size=13, color=CREAM)


def slide_closing():
    sl = prs.slides.add_slide(BLANK)
    fill_bg(sl, DARK)
    add_rect(sl, 0, 0, Inches(0.5), H, RED)

    add_text(sl, "D", Inches(0.8), Inches(0.6), Inches(1.5), Inches(1.5),
             size=80, bold=True, color=RED)

    add_text(sl, "RedBoard", Inches(0.8), Inches(1.8), Inches(10), Inches(1.0),
             size=52, bold=True, color=WHITE)

    add_text(sl, "One platform — built for the Big Red.",
             Inches(0.8), Inches(2.85), Inches(10), Inches(0.6),
             size=20, color=CREAM, italic=True)

    metrics = [
        ("5 Features", "Film · Stats · Wellness · Playbook · Feed"),
        ("2 Roles",    "Admin (coaching staff) · Player (roster)"),
        ("0 Servers",  "Fully serverless — Supabase + static hosting"),
        ("RLS",        "Every row protected — no client-side security holes"),
    ]

    for i, (key, val) in enumerate(metrics):
        y = Inches(3.7) + i * Inches(0.72)
        add_rect(sl, Inches(0.8), y, Inches(1.9), Inches(0.52), RED)
        add_text(sl, key, Inches(0.88), y + Inches(0.08), Inches(1.8),
                 Inches(0.42), size=14, bold=True, color=WHITE)
        add_text(sl, val, Inches(2.9), y + Inches(0.1), Inches(9.5),
                 Inches(0.42), size=14, color=CREAM)

    add_rect(sl, 0, H - Inches(0.5), W, Inches(0.5), RED)
    add_text(sl, "CONFIDENTIAL — INTERNAL USE ONLY",
             Inches(0.2), H - Inches(0.42), W - Inches(0.4), Inches(0.35),
             size=11, color=WHITE, align=PP_ALIGN.CENTER)


# ── Build all slides ───────────────────────────────────────────────────────────

slide_cover()

slide_problem()

slide_solution()

slide_stack()

slide_architecture()

# ── Login ──────────────────────────────────────────────────────────────────────
slide_feature_mobile(
    "Authentication",
    "Secure email + password login\nroutes each user to their role.",
    [
        "Supabase Auth issues a JWT on sign-in",
        "AuthContext reads role from the profiles table",
        "Admins land on /admin · Players land on /app/feed",
        "AuthGuard + AdminGuard enforce routing — no URL bypass",
    ],
    "00_login.png",
)

# ── Player Feed ────────────────────────────────────────────────────────────────
slide_feature_mobile(
    "Player Feed",
    "Every piece of coach content\nin one chronological timeline.",
    [
        "Time-based greeting with the player's first name",
        "Wellness nudge card when a check-in is due",
        "Stat snapshot — latest week vs prior, latest game vs career avg",
        "Unified feed: film clips, stat uploads, playbook additions",
        "Real-time updates via Supabase Realtime — no page refresh",
    ],
    "01_feed.png",
)

# ── Film ───────────────────────────────────────────────────────────────────────
slide_feature_mobile(
    "Film — Player View",
    "Hudl clips and video files\ndelivered directly to players.",
    [
        "Coaches target Team or Individual recipients",
        '"Just for you" badge on personal clips',
        "Unread red dot badge clears on page visit",
        "Hudl embedded inline · uploaded files via signed URL",
    ],
    "02_film.png",
)

# ── Stats ──────────────────────────────────────────────────────────────────────
slide_feature_mobile(
    "Stats — Player View",
    "Personal performance trends,\nhero stats, and coach goals.",
    [
        "Practice stats grouped by week with Δ vs prior week",
        "Game stats per game with Δ vs career average",
        "Four hero stats: Points · Rebounds · Assists · Steals",
        "Coach-set goals shown with progress bar",
        "All standard stats in expandable table",
    ],
    "03_stats.png",
)

# ── Wellness ───────────────────────────────────────────────────────────────────
slide_feature_mobile(
    "Wellness — Player Check-In",
    "One form per day.\nThree question types.",
    [
        "Rating (1–5 stars) · Yes/No toggle · Free-text response",
        "Active form shown only if player hasn't submitted today",
        "After submission: read-only review of answers",
        "Wellness nudge on Feed disappears on submission",
    ],
    "04_wellness.png",
)

# ── Playbook ───────────────────────────────────────────────────────────────────
slide_feature_mobile(
    "Playbook — Player View",
    "Virtual binder — browse folders,\ntap to open full-screen.",
    [
        "Folders expand to show file thumbnails",
        "Full-screen lightbox viewer for images",
        "PDFs render in-browser via iframe",
        "Coach controls folder and file sort order",
    ],
    "05_playbook.png",
)

# ── Profile ────────────────────────────────────────────────────────────────────
slide_feature_mobile(
    "Player Profile",
    "Players manage their own\nidentity within the platform.",
    [
        "Upload a profile photo (stored in Supabase Storage)",
        "Set jersey number, position, class year",
        "Toggle dim (dark) mode — persisted to localStorage",
        "Sign out",
    ],
    "06_profile.png",
)

# ── Admin section heading ──────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
fill_bg(sl, RED)
add_text(sl, "COACHING STAFF VIEW", Inches(1), Inches(2.5), Inches(11), Inches(0.6),
         size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Admin Features", Inches(1), Inches(3.1), Inches(11), Inches(1.4),
         size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Every feature has a coach-facing management interface.",
         Inches(1), Inches(4.5), Inches(11), Inches(0.6),
         size=20, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER, italic=True)

# ── Admin Dashboard ────────────────────────────────────────────────────────────
slide_feature_desktop(
    "Admin Dashboard",
    "Roster management and platform overview.",
    ["Create & manage user accounts", "Assign roles (admin / player)", "View profile completion"],
    "07_admin_dashboard.png",
)

# ── Admin Film ─────────────────────────────────────────────────────────────────
slide_feature_desktop(
    "Admin — Film",
    "Post Hudl links or upload video files to the whole team or individual players.",
    ["Paste Hudl URL or upload file", "Team or Individual visibility", "Select recipients from roster"],
    "08_admin_film.png",
)

# ── Admin Stats ────────────────────────────────────────────────────────────────
slide_feature_desktop(
    "Admin — Stats",
    "Upload CSV/XLSX → fuzzy-match roster → publish when ready.",
    ["Client-side parsing via SheetJS", "Auto column detection + manual override", "Draft → Publish flow  ·  Player annotations  ·  Goal setting"],
    "09_admin_stats.png",
)

# ── Admin Wellness ─────────────────────────────────────────────────────────────
slide_feature_desktop(
    "Admin — Wellness",
    "Build forms, schedule them in advance, and see today's responses at a glance.",
    ["Drag-and-drop form builder (Rating / Yes-No / Text)", "Calendar scheduling — auto-activate at midnight", "Today's dashboard: per-player status + team aggregates"],
    "10_admin_wellness.png",
)

# ── Admin Playbook ─────────────────────────────────────────────────────────────
slide_feature_desktop(
    "Admin — Playbook",
    "Create folders, upload files, drag to reorder.",
    ["Named folders with sort order", "Upload images or PDFs to any folder", "Drag-to-reorder within folders"],
    "11_admin_playbook.png",
)

slide_security()

slide_closing()

# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
